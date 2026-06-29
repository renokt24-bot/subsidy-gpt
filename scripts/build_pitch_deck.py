# -*- coding: utf-8 -*-
"""
USTORY next — Seed Pitch Deck generator.

Produces a 16:9 (1920x1080) fully-editable, Google-Slides-compatible .pptx.
Design system: White bg / Ink #111111 / Accent #2563EB / Gray #6B7280,
Inter (latin) + Noto Sans JP (japanese). Lucide icons rendered to PNG.

Run:  python3 scripts/build_pitch_deck.py
Out:  docs/USTORY_next_Pitch_Deck.pptx
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import cairosvg

# ----------------------------------------------------------------------------
# Paths
# ----------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
SCRATCH = "/tmp/claude-0/-home-user-subsidy-gpt/864517de-9424-5f97-ae9c-d5e57875190a/scratchpad"
SVG_DIR = os.path.join(SCRATCH, "icons_svg")
PNG_DIR = os.path.join(SCRATCH, "icons_png")
OUT = os.path.join(ROOT, "docs", "USTORY_next_Pitch_Deck.pptx")
os.makedirs(PNG_DIR, exist_ok=True)

# ----------------------------------------------------------------------------
# Design tokens
# ----------------------------------------------------------------------------
def C(hexstr): return RGBColor.from_string(hexstr)

WHITE   = "FFFFFF"
INK     = "111111"
INK2    = "374151"
ACCENT  = "2563EB"
ACCENT7 = "1D4ED8"
ACCENT_TINT  = "EEF3FE"
ACCENT_TINT2 = "DBE6FD"
GRAY    = "6B7280"
GRAY2   = "9CA3AF"
LINE    = "E5E7EB"
SOFT    = "F7F8FA"
SOFT2   = "F1F3F6"
GREEN   = "059669"
GREEN_TINT = "E7F6F0"
AMBER   = "D97706"
ROSE    = "E11D48"
ROSE_TINT  = "FCEAEE"

JP = "Noto Sans JP"
EN = "Inter"

# 1920x1080 canvas -> EMU (12192000 x 6858000), 6350 EMU/px
EMU_PX = 6350
def px(v): return Emu(int(round(v * EMU_PX)))

MARGIN = 132
CW = 1920 - 2 * MARGIN     # content width

L, Cn, R = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT
T, M, B = MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM

# ----------------------------------------------------------------------------
# Icon rendering (Lucide -> PNG, recolored)
# ----------------------------------------------------------------------------
_icon_cache = {}
def icon(name, color=INK, size=512, stroke=1.85):
    key = (name, color, size, stroke)
    if key in _icon_cache:
        return _icon_cache[key]
    out = os.path.join(PNG_DIR, f"{name}_{color}_{size}_{int(stroke*100)}.png")
    if not os.path.exists(out):
        with open(os.path.join(SVG_DIR, f"{name}.svg"), "r", encoding="utf-8") as f:
            svg = f.read()
        svg = svg.replace("currentColor", f"#{color}")
        svg = svg.replace('stroke-width="2"', f'stroke-width="{stroke}"')
        cairosvg.svg2png(bytestring=svg.encode("utf-8"),
                         write_to=out, output_width=size, output_height=size)
    _icon_cache[key] = out
    return out

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def _set_run_font(run, latin, ea):
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", latin if tag != "a:ea" else ea)

def tb(slide, x, y, w, h, content, size=20, color=INK, bold=False,
       align=L, anchor=T, latin=EN, ea=JP, line=1.16, space_after=0,
       spc=None, wrap=True):
    """Single-style text box. `content` may contain \\n for multiple lines."""
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = content.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        if space_after:
            p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = C(color)
        _set_run_font(run, latin, ea)
        if spc is not None:
            run._r.get_or_add_rPr().set("spc", str(int(spc * 100)))
    return box

def rect(slide, x, y, w, h, fill=WHITE, line_c=None, line_w=1.0,
         radius=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line_c is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line_c); sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if sp.has_text_frame:
        sp.text_frame.word_wrap = True
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(sp.text_frame, m, 0)
    return sp

def card(slide, x, y, w, h, fill=WHITE, line_c=LINE, radius=0.045, shadow=True):
    return rect(slide, x, y, w, h, fill=fill, line_c=line_c, line_w=1.0,
                radius=radius, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=shadow)

def _soft_shadow(sp):
    spPr = sp._element.spPr
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': '90000', 'dist': '38000',
                         'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '0F172A'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '9000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)

def line(slide, x1, y1, x2, y2, color=LINE, w=1.0, dash=None, cap_round=True):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    px(x1), px(y1), px(x2), px(y2))
    cn.line.color.rgb = C(color); cn.line.width = Pt(w)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if cap_round:
        ln.set('cap', 'rnd')
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    return cn

def add_icon(slide, name, x, y, size, color=INK, stroke=1.85):
    return slide.shapes.add_picture(icon(name, color, stroke=stroke),
                                    px(x), px(y), px(size), px(size))

def icon_tile(slide, x, y, size, name, tile=ACCENT_TINT, ic=ACCENT,
              radius=0.26, icon_ratio=0.52, stroke=1.9, line_c=None):
    rect(slide, x, y, size, size, fill=tile, line_c=line_c,
         radius=radius, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False)
    isz = size * icon_ratio
    add_icon(slide, name, x + (size - isz) / 2, y + (size - isz) / 2,
             isz, ic, stroke=stroke)

def pill(slide, x, y, w, h, text, fill=ACCENT_TINT, color=ACCENT7,
         size=13, bold=True, line_c=None, latin=EN):
    rect(slide, x, y, w, h, fill=fill, line_c=line_c, radius=0.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False)
    tb(slide, x, y - 1, w, h, text, size=size, color=color, bold=bold,
       align=Cn, anchor=M, latin=latin)

# ----------------------------------------------------------------------------
# Slide chrome
# ----------------------------------------------------------------------------
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(WHITE)
    return s

def header(slide, kicker, title, subtitle=None, title_size=38):
    # accent kicker with a small leading tick
    rect(slide, MARGIN, 95, 26, 4, fill=ACCENT, radius=0.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(slide, MARGIN + 38, 86, CW - 38, 24, kicker, size=13, color=ACCENT7,
       bold=True, spc=2.4, anchor=M)
    tb(slide, MARGIN, 122, CW, title_size * 1.5, title, size=title_size,
       color=INK, bold=True, line=1.05)
    if subtitle:
        tb(slide, MARGIN, 214, CW, 36, subtitle,
           size=19, color=GRAY, line=1.2)

def footer(slide, n, total=18):
    tb(slide, MARGIN, 1019, 400, 22, "USTORY next", size=11, color=GRAY2,
       bold=True, anchor=M)
    tb(slide, 1920 - MARGIN - 300, 1019, 300, 22, f"{n:02d} / {total}",
       size=11, color=GRAY2, align=R, anchor=M)

def source(slide, text, y=986, x=None):
    if x is None: x = MARGIN
    tb(slide, x, y, CW, 20, text, size=10.5, color=GRAY2, line=1.15)

# ----------------------------------------------------------------------------
# Composite components
# ----------------------------------------------------------------------------
def stat_card(slide, x, y, w, h, value, label, sub=None, icon_name=None,
              accent=ACCENT, vsize=42, fill=WHITE, line_c=LINE):
    card(slide, x, y, w, h, fill=fill, line_c=line_c)
    pad = 26
    iy = y + pad
    if icon_name:
        icon_tile(slide, x + pad, iy, 46, icon_name)
        iy += 64
    tb(slide, x + pad, iy if icon_name else y + pad + 4, w - 2 * pad, vsize * 1.4,
       value, size=vsize, color=INK, bold=True, line=1.0)
    tb(slide, x + pad, iy + (vsize * 1.18 if True else 0) + (0 if icon_name else 0),
       w - 2 * pad, 26, label, size=14.5, color=INK2, bold=True, line=1.15)
    if sub:
        tb(slide, x + pad, y + h - pad - 20, w - 2 * pad, 22, sub,
           size=12, color=GRAY, line=1.15)

def feature_row(slide, x, y, w, icon_name, title, desc, ic=ACCENT,
                tile=ACCENT_TINT, tsize=16, dsize=13):
    icon_tile(slide, x, y, 40, icon_name, tile=tile, ic=ic)
    tb(slide, x + 56, y - 2, w - 56, 24, title, size=tsize, color=INK, bold=True)
    tb(slide, x + 56, y + 24, w - 56, 40, desc, size=dsize, color=GRAY, line=1.25)

def node(slide, x, y, w, h, icon_name, title, desc, accent=ACCENT,
         tile=ACCENT_TINT, fill=WHITE, line_c=LINE, title_size=18):
    card(slide, x, y, w, h, fill=fill, line_c=line_c)
    icon_tile(slide, x + 24, y + 24, 50, icon_name, tile=tile, ic=accent)
    tb(slide, x + 24, y + 88, w - 48, 26, title, size=title_size, color=INK, bold=True)
    tb(slide, x + 24, y + 88 + title_size + 8, w - 48, h - 120, desc,
       size=13, color=GRAY, line=1.35)

# ----------------------------------------------------------------------------
# SLIDE 1 — Cover
# ----------------------------------------------------------------------------
def s01_cover(prs):
    s = new_slide(prs)
    # subtle right-side node graphic panel
    panelx = 1180
    rect(s, panelx, 0, 1920 - panelx, 1080, fill=SOFT)
    line(s, panelx, 0, panelx, 1080, color=LINE, w=1.0)
    # node triangle (学校/学生/企業 around USTORY hub)
    cx, cy = 1545, 540
    pts = {"学校": (cx, cy - 210), "学生": (cx - 195, cy + 150),
           "企業": (cx + 195, cy + 150)}
    for (lx, ly) in pts.values():
        line(s, cx, cy, lx, ly, color=ACCENT_TINT2, w=2.2)
    icons = {"学校": "school", "学生": "graduation-cap", "企業": "building-2"}
    for lbl, (lx, ly) in pts.items():
        rect(s, lx - 52, ly - 52, 104, 104, fill=WHITE, line_c=LINE,
             radius=0.22, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        add_icon(s, icons[lbl], lx - 23, ly - 28, 46, ACCENT, stroke=1.8)
        tb(s, lx - 52, ly + 30, 104, 20, lbl, size=12, color=INK2, bold=True, align=Cn)
    # hub
    rect(s, cx - 58, cy - 58, 116, 116, fill=ACCENT, radius=0.24,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    tb(s, cx - 58, cy - 58, 116, 116, "U", size=52, color=WHITE, bold=True,
       align=Cn, anchor=M, latin=EN)
    # left content
    rect(s, MARGIN, 360, 30, 30, fill=ACCENT, radius=0.28,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, MARGIN + 44, 360, 700, 30, "USTORY next", size=17, color=INK, bold=True)
    tb(s, MARGIN, 432, 1000, 30, "BUILD JAPAN'S", size=15, color=ACCENT7,
       bold=True, spc=3.0)
    tb(s, MARGIN, 466, 1040, 280, "Career\nInfrastructure.", size=62,
       color=INK, bold=True, line=1.06)
    tb(s, MARGIN, 742, 1040, 40,
       "才能が、学歴や履歴書に埋もれない社会をつくる。", size=22, color=INK2, line=1.3)
    line(s, MARGIN, 806, MARGIN + 96, 806, color=ACCENT, w=2.5)
    tb(s, MARGIN, 828, 1040, 30,
       "専門学校を起点に、学校・学生・企業をつなぐキャリアインフラ。", size=15,
       color=GRAY, line=1.5)
    # bottom meta
    tb(s, MARGIN, 980, 600, 24, "Seed Round  —  ¥70,000,000", size=14,
       color=INK, bold=True)
    tb(s, 1920 - MARGIN - 300 - (1920 - panelx), 980, 300, 24, "Pitch Deck 2026",
       size=13, color=GRAY, align=R)

# ----------------------------------------------------------------------------
# SLIDE 2 — Executive Summary
# ----------------------------------------------------------------------------
def s02_exec(prs):
    s = new_slide(prs)
    header(s, "EXECUTIVE SUMMARY", "専門学校発、日本のキャリアインフラ。",
           "学校・学生・企業をつなぎ、スキルで採用が決まる市場をつくる。")
    cards = [
        ("layers", "事業", "キャリアインフラ",
         "専門学校を起点に学校・学生・企業を接続するSaaS＋マッチング基盤。"),
        ("trending-up", "市場", "SAM ¥2,851億",
         "新卒採用支援＋HR Techクラウド。年20〜30%成長の追い風。"),
        ("route", "優位", "学校営業モデル",
         "AIではなくDistribution。高校市場で実証済みの導入運用網。"),
        ("rocket", "現在地", "MVP完成・7月始動",
         "5校／企業10社が始動。2026年7月に正式リリース。"),
    ]
    n = len(cards); gap = 26
    w = (CW - gap * (n - 1)) / n; y = 360; h = 360
    for i, (ic, tag, val, desc) in enumerate(cards):
        x = MARGIN + i * (w + gap)
        card(s, x, y, w, h)
        icon_tile(s, x + 28, y + 30, 52, ic)
        tb(s, x + 28, y + 100, w - 56, 20, tag, size=12, color=ACCENT7,
           bold=True, spc=1.5)
        tb(s, x + 28, y + 128, w - 56, 64, val, size=21, color=INK, bold=True, line=1.1)
        line(s, x + 28, y + 200, x + w - 28, y + 200, color=LINE, w=1.0)
        tb(s, x + 28, y + 220, w - 56, h - 240, desc, size=13.5, color=GRAY, line=1.45)
    # thesis strip
    by = 752
    rect(s, MARGIN, by, CW, 78, fill=INK, radius=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_icon(s, "quote", 0, 0, 1, ACCENT) if False else None
    tb(s, MARGIN + 34, by, CW - 380, 78,
       "“ 専門スキルを持つ55万人の才能が、履歴書ではなくポートフォリオで評価される。”",
       size=17, color=WHITE, bold=True, anchor=M, line=1.3)
    pill(s, 1920 - MARGIN - 300, by + 21, 272, 36,
         "Skill  ＞  Resume", fill=ACCENT, color=WHITE, size=15, latin=EN)
    footer(s, 2)

# ----------------------------------------------------------------------------
# SLIDE 3 — Investment Highlights
# ----------------------------------------------------------------------------
def s03_highlights(prs):
    s = new_slide(prs)
    header(s, "INVESTMENT HIGHLIGHTS", "なぜ、いま投資すべきか。",
           "巨大市場 × 独自Distribution × ネットワーク効果 × 資本効率。")
    items = [
        ("trending-up", "巨大かつ成長市場", "SAM ¥2,851億",
         "新卒採用支援とHR Techクラウドが年20〜30%で拡大。採用単価も高騰。"),
        ("route", "模倣困難なDistribution", "学校営業モデル",
         "高校市場で学校営業→導入→CS→運用を構築済み。専門学校へ横展開。"),
        ("repeat", "強いネットワーク効果", "三者データの蓄積",
         "学校・学生・企業が増えるほど価値が上がり、作品/スキルデータが堀になる。"),
        ("badge-check", "実証済みの初期実績", "MVP・5校・10社",
         "プロダクト完成。専門学校5校・導入可能企業10社で2026年7月始動。"),
        ("gauge", "高い資本効率", "LTV/CAC 6.6x",
         "学校チャネルでCACを抑制。Payback 5ヶ月、粗利80%超の構造。"),
        ("users", "実行力あるチーム", "高校市場の構築者",
         "学校営業の現場を作ったFounder陣。CEO・エンジニアの三位一体。"),
    ]
    cols, rows = 3, 2; gap = 26
    w = (CW - gap * (cols - 1)) / cols
    h = 230; y0 = 350; vgap = 24
    for i, (ic, title, val, desc) in enumerate(items):
        r, c = divmod(i, cols)
        x = MARGIN + c * (w + gap)
        y = y0 + r * (h + vgap)
        card(s, x, y, w, h)
        rect(s, x, y, 5, h, fill=ACCENT, radius=0.0)
        icon_tile(s, x + 26, y + 26, 46, ic)
        tb(s, x + 86, y + 28, w - 110, 20, title, size=13, color=GRAY, bold=True)
        tb(s, x + 86, y + 50, w - 110, 32, val, size=22, color=INK, bold=True)
        tb(s, x + 26, y + 118, w - 52, h - 130, desc, size=13, color=GRAY, line=1.4)
    footer(s, 3)

# ----------------------------------------------------------------------------
# SLIDE 4 — Why Now
# ----------------------------------------------------------------------------
def s04_whynow(prs):
    s = new_slide(prs)
    header(s, "WHY NOW", "3つの追い風が、いま重なる。",
           "採用コストの高騰、HR Techの急成長、そして専門学校DXの空白。")
    forces = [
        ("banknote", "採用単価の高騰", "¥45.0万 → ¥56.8万",
         "新卒採用単価は1年で+26%。母集団形成型の採用は限界に。",
         "出典：マイナビ 2024年卒 企業新卒内定状況調査", "+26% YoY"),
        ("zap", "HR Techクラウドの急成長", "¥1,385億 → ¥3,200億",
         "2022→2027でCAGR +31.8%。採用はSaaS×データへ移行中。",
         "出典：ミック経済研究所 HRTechクラウド市場", "CAGR +31.8%"),
        ("school", "専門学校DXの空白", "2,676校 が未DX",
         "求人・就職管理はアナログ中心。スキル人材の流通網が存在しない。",
         "出典：文部科学省 令和6年度 学校基本調査", "0 → 1 の市場"),
    ]
    n = 3; gap = 28
    w = (CW - gap * (n - 1)) / n; y = 356; h = 470
    for i, (ic, title, val, desc, src, badge) in enumerate(forces):
        x = MARGIN + i * (w + gap)
        card(s, x, y, w, h)
        icon_tile(s, x + 30, y + 32, 58, ic, radius=0.28)
        pill(s, x + w - 30 - 150, y + 40, 150, 34, badge, fill=ACCENT_TINT,
             color=ACCENT7, size=12.5)
        tb(s, x + 30, y + 116, w - 60, 24, title, size=15, color=GRAY, bold=True)
        tb(s, x + 30, y + 150, w - 60, 60, val, size=26, color=INK, bold=True, line=1.05)
        line(s, x + 30, y + 246, x + w - 30, y + 246, color=LINE, w=1.0)
        tb(s, x + 30, y + 270, w - 60, 120, desc, size=14, color=INK2, line=1.5)
        tb(s, x + 30, y + h - 52, w - 60, 36, src, size=10.5, color=GRAY2, line=1.3)
        # up arrow accent
        add_icon(s, "trending-up", x + w - 30 - 26, y + h - 50, 24, ACCENT, stroke=2.0)
    footer(s, 4)

# ----------------------------------------------------------------------------
# SLIDE 5 — Problem
# ----------------------------------------------------------------------------
def s05_problem(prs):
    s = new_slide(prs)
    header(s, "PROBLEM", "才能が、履歴書に埋もれている。",
           "学校・学生・企業の三者が、それぞれの断絶を抱えている。")
    # central broken-link triangle
    gx, gy = 1430, 600
    coords = {"学生": (gx, gy - 165), "学校": (gx - 175, gy + 135),
              "企業": (gx + 175, gy + 135)}
    # broken (dashed) links with X
    keys = list(coords.keys())
    for a in range(len(keys)):
        for b2 in range(a + 1, len(keys)):
            (x1, y1), (x2, y2) = coords[keys[a]], coords[keys[b2]]
            line(s, x1, y1, x2, y2, color=ROSE, w=1.8, dash="dash")
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            rect(s, mx - 15, my - 15, 30, 30, fill=WHITE, line_c=ROSE,
                 line_w=1.4, radius=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_icon(s, "x", mx - 8, my - 8, 16, ROSE, stroke=2.4)
    micons = {"学生": "graduation-cap", "学校": "school", "企業": "building-2"}
    for lbl, (x1, y1) in coords.items():
        rect(s, x1 - 58, y1 - 58, 116, 116, fill=WHITE, line_c=LINE,
             radius=0.22, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        add_icon(s, micons[lbl], x1 - 24, y1 - 30, 48, INK, stroke=1.7)
        tb(s, x1 - 58, y1 + 28, 116, 20, lbl, size=13, color=INK, bold=True, align=Cn)
    # left: three pain points
    pains = [
        ("graduation-cap", "学生", "スキルが、伝わらない。",
         "作品・実績があっても、履歴書とエントリーシートでしか評価されない。"),
        ("building-2", "企業", "専門人材に、出会えない。",
         "即戦力スキルを持つ学生が見えず、採用単価だけが上がり続ける。"),
        ("school", "学校", "就職支援が、アナログ。",
         "求人・説明会・学生管理が紙とExcel。企業との接点を作れない。"),
    ]
    y = 350; ph = 150; gap = 22; lw = 1080
    for i, (ic, who, t, d) in enumerate(pains):
        yy = y + i * (ph + gap)
        card(s, MARGIN, yy, lw, ph)
        icon_tile(s, MARGIN + 26, yy + 30, 56, ic, tile=ROSE_TINT, ic=ROSE)
        tb(s, MARGIN + 104, yy + 30, 120, 22, who, size=13, color=ROSE, bold=True)
        tb(s, MARGIN + 104, yy + 54, lw - 130, 34, t, size=21, color=INK, bold=True)
        tb(s, MARGIN + 104, yy + 96, lw - 130, 40, d, size=13.5, color=GRAY, line=1.35)
    footer(s, 5)

# ----------------------------------------------------------------------------
# SLIDE 6 — Existing Solutions
# ----------------------------------------------------------------------------
def s06_existing(prs):
    s = new_slide(prs)
    header(s, "EXISTING SOLUTIONS", "既存サービスは、三者を分断している。",
           "それぞれが一辺だけを担い、学校・学生・企業を1つに繋ぐ主体がいない。")
    cats = [
        ("search", "求人媒体", "マイナビ / リクナビ",
         ["学生", "企業"], "大学中心の母集団形成。専門スキルは可視化されない。"),
        ("send", "スカウト媒体", "OfferBox / dodaキャンパス",
         ["学生", "企業"], "プロフィール起点の逆求人。学校との連携は持たない。"),
        ("folder-search", "学校管理SaaS", "Campus Plan / Sakura",
         ["学校"], "校務を効率化するが、企業・採用とは接続しない。"),
        ("image", "ポートフォリオ", "foriio / Behance",
         ["学生"], "作品を見せる場。採用・スカウトには接続しない。"),
    ]
    n = 4; gap = 24
    w = (CW - gap * (n - 1)) / n; y = 352; h = 392
    chips = ["学校", "学生", "企業"]
    for i, (ic, name, ex, covers, desc) in enumerate(cats):
        x = MARGIN + i * (w + gap)
        card(s, x, y, w, h)
        icon_tile(s, x + 26, y + 28, 50, ic, tile=SOFT2, ic=INK2)
        tb(s, x + 26, y + 92, w - 52, 24, name, size=17, color=INK, bold=True)
        tb(s, x + 26, y + 120, w - 52, 20, ex, size=12, color=GRAY)
        # coverage chips (filled = covered)
        cy = y + 168
        tb(s, x + 26, cy - 26, w - 52, 18, "カバー範囲", size=11, color=GRAY2, bold=True, spc=1.0)
        cw = (w - 52 - 16) / 3
        for j, ch in enumerate(chips):
            cx = x + 26 + j * (cw + 8)
            on = ch in covers
            rect(s, cx, cy, cw, 38, fill=ACCENT_TINT if on else SOFT,
                 line_c=None if on else LINE, radius=0.30,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tb(s, cx, cy - 1, cw, 38, ch, size=12.5,
               color=ACCENT7 if on else GRAY2, bold=on, align=Cn, anchor=M)
        line(s, x + 26, y + 252, x + w - 26, y + 252, color=LINE, w=1.0)
        tb(s, x + 26, y + 270, w - 52, 100, desc, size=12.5, color=GRAY, line=1.4)
    # bottom takeaway
    by = 786
    rect(s, MARGIN, by, CW, 64, fill=ACCENT_TINT, radius=0.14,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_icon(s, "scan-search", MARGIN + 24, by + 17, 30, ACCENT7, stroke=2.0)
    tb(s, MARGIN + 70, by, CW - 100, 64,
       "市場は断片化。三者を同時につなぐ「キャリアインフラ」の空白が残っている。",
       size=16, color=ACCENT7, bold=True, anchor=M)
    footer(s, 6)

# ----------------------------------------------------------------------------
# SLIDE 7 — Solution
# ----------------------------------------------------------------------------
def s07_solution(prs):
    s = new_slide(prs)
    header(s, "SOLUTION", "USTORY next が、三者を1つにつなぐ。",
           "学校を起点に、学生のスキルと企業の採用を1つの基盤で循環させる。")
    detail = {
        "学生": ("graduation-cap", ["作品・動画・スキル", "ポートフォリオ公開", "企業からスカウト"]),
        "学校": ("school", ["求人・説明会を管理", "学生ポートフォリオ", "就職実績を可視化"]),
        "企業": ("building-2", ["スキル検索・AI", "ダイレクトスカウト", "採用管理まで一気通貫"]),
    }
    cx, cy = 960, 660
    hw, hh = 240, 150
    boxw, boxh = 384, 196
    # placements
    stu = (cx - boxw / 2, 318)                 # top center
    sch = (MARGIN, 560)                        # left
    com = (1920 - MARGIN - boxw, 560)          # right
    # connectors hub <-> cards
    line(s, cx, cy - hh / 2, cx, stu[1] + boxh, color=ACCENT_TINT2, w=3.0)
    line(s, cx - hw / 2, cy, sch[0] + boxw, sch[1] + boxh / 2, color=ACCENT_TINT2, w=3.0)
    line(s, cx + hw / 2, cy, com[0], com[1] + boxh / 2, color=ACCENT_TINT2, w=3.0)
    # hub
    rect(s, cx - hw / 2, cy - hh / 2, hw, hh, fill=ACCENT, radius=0.14,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    add_icon(s, "repeat", cx - 22, cy - 52, 44, WHITE, stroke=1.9)
    tb(s, cx - hw / 2, cy + 2, hw, 28, "USTORY next", size=18, color=WHITE,
       bold=True, align=Cn)
    tb(s, cx - hw / 2, cy + 32, hw, 24, "キャリアインフラ", size=12.5,
       color="C9DBFE", align=Cn)
    # info cards
    for lbl, (bx, by) in {"学生": stu, "学校": sch, "企業": com}.items():
        card(s, bx, by, boxw, boxh)
        icon_tile(s, bx + 24, by + 24, 46, detail[lbl][0])
        tb(s, bx + 84, by + 24, boxw - 104, 46, lbl, size=18, color=INK,
           bold=True, anchor=M)
        line(s, bx + 24, by + 86, bx + boxw - 24, by + 86, color=LINE, w=1.0)
        for k, feat in enumerate(detail[lbl][1]):
            fy = by + 104 + k * 30
            rect(s, bx + 26, fy + 8, 7, 7, fill=ACCENT, radius=0.5,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            tb(s, bx + 44, fy, boxw - 68, 24, feat, size=13.5, color=INK2, anchor=M)
    footer(s, 7)

# ----------------------------------------------------------------------------
# SLIDE 8 — Product
# ----------------------------------------------------------------------------
def s08_product(prs):
    s = new_slide(prs)
    header(s, "PRODUCT", "3つのプロダクトが、1つにつながる。",
           "学生・企業・学校、それぞれの体験を1つのクラウド上で。")
    cols = [
        ("graduation-cap", "学生", "Portfolio",
         [("image", "作品・動画を掲載"), ("badge-check", "スキルを構造化"),
          ("send", "企業からスカウト受信")]),
        ("building-2", "企業", "Hiring Cloud",
         [("search", "スキル検索"), ("sparkles", "AIマッチング"),
          ("clipboard-list", "スカウト〜採用管理")]),
        ("school", "学校", "School Cloud",
         [("briefcase", "求人管理"), ("megaphone", "説明会管理"),
          ("users", "学生管理")]),
    ]
    n = 3; gap = 28
    w = (CW - gap * (n - 1)) / n; y = 352; h = 484
    for i, (ic, who, en, feats) in enumerate(cols):
        x = MARGIN + i * (w + gap)
        card(s, x, y, w, h)
        # header band
        icon_tile(s, x + 26, y + 26, 50, ic)
        tb(s, x + 88, y + 28, w - 110, 22, who, size=13, color=ACCENT7, bold=True)
        tb(s, x + 88, y + 50, w - 110, 28, en, size=20, color=INK, bold=True)
        # device placeholder (image placeholder)
        dy = y + 100
        rect(s, x + 26, dy, w - 52, 168, fill=SOFT, line_c=LINE, radius=0.05,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x + 26, dy, w - 52, 30, fill=SOFT2, radius=0.0)
        for d in range(3):
            rect(s, x + 42 + d * 16, dy + 12, 7, 7, fill=GRAY2, radius=0.5,
                 shape=MSO_SHAPE.OVAL)
        add_icon(s, "monitor-smartphone", x + w / 2 - 22, dy + 70, 44, GRAY2, stroke=1.6)
        tb(s, x + 26, dy + 120, w - 52, 20, "UI placeholder", size=11,
           color=GRAY2, align=Cn)
        # features
        fy0 = dy + 192
        for k, (fic, ft) in enumerate(feats):
            fy = fy0 + k * 42
            icon_tile(s, x + 26, fy, 30, fic, radius=0.30, icon_ratio=0.56)
            tb(s, x + 68, fy + 1, w - 94, 28, ft, size=14.5, color=INK2,
               bold=True, anchor=M)
    footer(s, 8)

# ----------------------------------------------------------------------------
# SLIDE 9 — Why USTORY (Differentiation)
# ----------------------------------------------------------------------------
def s09_why(prs):
    s = new_slide(prs)
    header(s, "WHY USTORY", "競争優位は、AIではない。学校営業モデル。",
           "プロダクトは模倣される。模倣できないのは、学校への流通網。")
    # left: contrast two boxes
    lw = 540; y = 352
    card(s, MARGIN, y, lw, 200, fill=SOFT, line_c=LINE)
    tb(s, MARGIN + 28, y + 24, lw - 56, 22, "他社の競争軸", size=13, color=GRAY, bold=True)
    tb(s, MARGIN + 28, y + 52, lw - 56, 40, "プロダクト勝負", size=26, color=INK2, bold=True)
    tb(s, MARGIN + 28, y + 108, lw - 56, 70,
       "機能・AI・UIで差別化 → すぐに追いつかれ、価格競争に陥る。",
       size=14, color=GRAY, line=1.5)
    card(s, MARGIN, y + 224, lw, 200, fill=ACCENT, line_c=None)
    tb(s, MARGIN + 28, y + 248, lw - 56, 22, "USTORYの競争軸", size=13,
       color="C9DBFE", bold=True)
    tb(s, MARGIN + 28, y + 276, lw - 56, 40, "Distribution勝負", size=26,
       color=WHITE, bold=True)
    tb(s, MARGIN + 28, y + 332, lw - 56, 70,
       "学校との信頼関係と運用網。後発でも入り込めない参入障壁になる。",
       size=14, color="E3ECFE", line=1.5)
    # right: the distribution engine pipeline (built in HS -> expand to vocational)
    rx = MARGIN + lw + 48; rw = CW - lw - 48
    card(s, rx, y, rw, 472)
    tb(s, rx + 32, y + 28, rw - 64, 24, "実証済みの学校営業エンジン", size=16,
       color=INK, bold=True)
    pill(s, rx + rw - 32 - 196, y + 26, 196, 34, "高校市場で構築済み",
         fill=GREEN_TINT, color=GREEN, size=12)
    steps = [("handshake", "学校営業"), ("badge-check", "導入"),
             ("headset", "CS"), ("settings", "運用")]
    sx = rx + 40; sy = y + 110; sw = (rw - 80 - 3 * 40) / 4
    for k, (ic, label) in enumerate(steps):
        bx = sx + k * (sw + 40)
        rect(s, bx, sy, sw, 116, fill=SOFT, line_c=LINE, radius=0.10,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        icon_tile(s, bx + sw / 2 - 25, sy + 16, 50, ic)
        tb(s, bx, sy + 74, sw, 26, label, size=14, color=INK, bold=True, align=Cn)
        if k < 3:
            ax = bx + sw + 6
            add_icon(s, "arrow-right", ax + 6, sy + 44, 26, ACCENT, stroke=2.0)
    # expansion arrow to vocational
    ey = y + 270
    line(s, rx + 40, ey, rx + rw - 40, ey, color=LINE, w=1.0)
    tb(s, rx + 40, ey + 18, rw - 80, 22, "そのまま専門学校市場へ横展開", size=14,
       color=GRAY, bold=True)
    # from-to bar
    fy = ey + 56
    rect(s, rx + 40, fy, (rw - 80) * 0.42, 70, fill=GREEN_TINT, radius=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, rx + 60, fy + 12, 260, 24, "高校市場", size=13, color=GREEN, bold=True)
    tb(s, rx + 60, fy + 36, 260, 24, "学校営業を構築・実証", size=13, color=INK2)
    add_icon(s, "arrow-right", rx + 40 + (rw - 80) * 0.44, fy + 20, 28, INK, stroke=2.2)
    bx2 = rx + 40 + (rw - 80) * 0.52
    rect(s, bx2, fy, (rw - 80) * 0.48, 70, fill=ACCENT, radius=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, bx2 + 22, fy + 12, 320, 24, "専門学校市場 — 2,676校", size=13,
       color=WHITE, bold=True)
    tb(s, bx2 + 22, fy + 36, 320, 24, "同じ営業網で一気に獲得", size=13, color="E3ECFE")
    footer(s, 9)

# ----------------------------------------------------------------------------
# SLIDE 10 — Market (TAM SAM SOM)
# ----------------------------------------------------------------------------
def s10_market(prs):
    s = new_slide(prs)
    header(s, "MARKET", "1兆円規模の市場に、最短距離で入る。",
           "人材ビジネス市場を頂点に、専門学校採用というビーコンから攻める。")
    # concentric circles on left
    ccx, ccy = 548, 588
    rings = [(384, ACCENT_TINT), (266, ACCENT_TINT2), (158, ACCENT)]
    for d, col in rings:
        rect(s, ccx - d / 2, ccy - d / 2, d, d, fill=col,
             shape=MSO_SHAPE.OVAL, line_c=None)
    tb(s, ccx - 80, ccy - 38, 160, 30, "¥140億", size=21, color=WHITE, bold=True, align=Cn)
    tb(s, ccx - 80, ccy - 8, 160, 20, "SOM", size=11.5, color="C9DBFE", align=Cn, bold=True)
    tb(s, ccx - 150, ccy + 50, 300, 22, "SAM · ¥2,851億", size=12.5, color=ACCENT7,
       bold=True, align=Cn)
    tb(s, ccx - 220, ccy + 108, 440, 22, "TAM · ¥10.2兆", size=12.5, color=GRAY,
       bold=True, align=Cn)
    # right: definitions
    rx = 1010; rw = 1920 - MARGIN - rx; y = 318
    defs = [
        ("TAM", "¥10.2兆", "日本の人材ビジネス市場",
         "人材派遣・紹介・再就職支援の総市場。", GRAY, SOFT),
        ("SAM", "¥2,851億", "新卒採用支援 ＋ HR Tech",
         "¥1,466億 ＋ ¥1,385億。年20〜30%で成長。", ACCENT7, ACCENT_TINT),
        ("SOM", "¥140億", "専門学校卒の採用市場",
         "専門卒採用 約¥1,400億の10%を5年で獲得。", WHITE, ACCENT),
    ]
    h = 132; gap = 18; vw = 308
    for i, (tag, val, t, d, tc, fill) in enumerate(defs):
        yy = y + i * (h + gap)
        is_accent = (fill == ACCENT)
        card(s, rx, yy, rw, h, fill=fill, line_c=None if is_accent else LINE, shadow=not is_accent)
        tb(s, rx + 30, yy + 26, vw - 36, 24, tag, size=14,
           color="C9DBFE" if is_accent else tc, bold=True, spc=1.5)
        tb(s, rx + 30, yy + 54, vw - 36, 50, val, size=28,
           color=WHITE if is_accent else INK, bold=True, line=1.0)
        line(s, rx + vw, yy + 28, rx + vw, yy + h - 28,
             color="3B7BF0" if is_accent else LINE, w=1.0)
        tb(s, rx + vw + 30, yy + 30, rw - vw - 58, 24, t, size=16,
           color=WHITE if is_accent else INK, bold=True)
        tb(s, rx + vw + 30, yy + 64, rw - vw - 58, 50, d, size=12.5,
           color="E3ECFE" if is_accent else GRAY, line=1.4)
    # building blocks strip
    by = y + 3 * (h + gap) + 30
    blocks = [("2,676校", "専門学校数"), ("558,255名", "専門学校生"),
              ("¥56.8万", "新卒採用単価"), ("+31.8%", "HR Tech CAGR")]
    bw = (CW - 3 * 16) / 4
    for i, (v, l) in enumerate(blocks):
        bx = MARGIN + i * (bw + 16)
        rect(s, bx, by, bw, 80, fill=SOFT, line_c=LINE, radius=0.12,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb(s, bx + 22, by + 14, bw - 44, 30, v, size=20, color=ACCENT7, bold=True, line=1.0)
        tb(s, bx + 22, by + 52, bw - 44, 20, l, size=11.5, color=GRAY)
    source(s, "出典：文部科学省 令和6年度学校基本調査／矢野経済研究所（新卒採用支援・人材ビジネス）／ミック経済研究所（HRTechクラウド）／マイナビ 2024年卒調査")
    footer(s, 10)

# ----------------------------------------------------------------------------
# SLIDE 11 — Business Model
# ----------------------------------------------------------------------------
def s11_model(prs):
    s = new_slide(prs)
    header(s, "BUSINESS MODEL", "収益は、4段階で積み上がる。",
           "企業SaaSから始め、スカウト・採用課金・AIへと単価を拡張する。")
    phases = [
        ("layers", "Phase 1", "企業SaaS", "月額課金",
         "企業の採用管理を月額SaaSで提供。安定したMRRの土台。", "¥3万/月〜"),
        ("send", "Phase 2", "スカウト", "従量課金",
         "学生プロフィールへのダイレクトスカウトを従量で提供。", "ARPA拡大"),
        ("badge-check", "Phase 3", "採用課金", "成果報酬",
         "採用決定時の成功報酬。高単価レイヤーを上積み。", "成功報酬"),
        ("sparkles", "Phase 4", "AI", "付加価値",
         "AIマッチング・スキル評価で単価と継続率をさらに向上。", "LTV最大化"),
    ]
    n = 4; gap = 24
    w = (CW - gap * (n - 1)) / n; y = 360; h = 372
    for i, (ic, ph, name, model, desc, tagv) in enumerate(phases):
        x = MARGIN + i * (w + gap)
        # rising step illusion: top offset decreases
        off = (3 - i) * 22
        card(s, x, y + off, w, h - off)
        rect(s, x, y + off, w, 6, fill=ACCENT, radius=0.0)
        icon_tile(s, x + 26, y + off + 26, 50, ic)
        pill(s, x + w - 26 - 96, y + off + 34, 96, 32, tagv, fill=ACCENT_TINT,
             color=ACCENT7, size=11.5)
        tb(s, x + 26, y + off + 92, w - 52, 20, ph, size=12, color=ACCENT7,
           bold=True, spc=1.5)
        tb(s, x + 26, y + off + 114, w - 52, 34, name, size=23, color=INK, bold=True)
        tb(s, x + 26, y + off + 150, w - 52, 20, model, size=12.5, color=GRAY, bold=True)
        line(s, x + 26, y + off + 184, x + w - 26, y + off + 184, color=LINE, w=1.0)
        tb(s, x + 26, y + off + 202, w - 52, h - off - 220, desc, size=13,
           color=GRAY, line=1.45)
    # base note
    by = 770
    rect(s, MARGIN, by, CW, 58, fill=SOFT, radius=0.14,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_icon(s, "trending-up", MARGIN + 24, by + 15, 28, ACCENT7, stroke=2.0)
    tb(s, MARGIN + 66, by, CW - 90, 58,
       "学校チャネルで企業を低CACで獲得 → SaaSで土台 → スカウト・採用課金でARPAを拡張する積み上げ型モデル。",
       size=14.5, color=INK2, bold=True, anchor=M)
    footer(s, 11)

# ----------------------------------------------------------------------------
# SLIDE 12 — Go To Market
# ----------------------------------------------------------------------------
def s12_gtm(prs):
    s = new_slide(prs)
    header(s, "GO TO MARKET", "学校を入口に、企業を面で取る。",
           "学校1校の導入が、数百人の学生と数十社の企業を連れてくる。")
    # left funnel
    fx = MARGIN; fw = 772; y = 352
    steps = [
        ("school", "専門学校を開拓", "学校営業モデルで導入", "50", "校", ACCENT),
        ("graduation-cap", "学生がオンボード", "ポートフォリオを公開", "数千", "名", ACCENT),
        ("building-2", "企業が無料で参加", "学生を検索・スカウト", "3,000", "社", ACCENT),
        ("credit-card", "企業が有料転換", "SaaS・スカウトで課金", "180", "社", ACCENT7),
    ]
    h = 96; gap = 24; vzone = 168
    for i, (ic, t, d, val, unit, col) in enumerate(steps):
        yy = y + i * (h + gap)
        ww = fw - i * 56
        bx = fx + (fw - ww) / 2
        card(s, bx, yy, ww, h, fill=ACCENT_TINT if i == 3 else WHITE,
             line_c=None if i == 3 else LINE)
        icon_tile(s, bx + 20, yy + 23, 50, ic)
        tb(s, bx + 86, yy + 24, ww - 96 - vzone, 26, t, size=16, color=INK, bold=True)
        tb(s, bx + 86, yy + 52, ww - 96 - vzone, 24, d, size=12.5, color=GRAY)
        # value zone, right
        vx = bx + ww - vzone
        tb(s, vx, yy, vzone - 44, h, val, size=27, color=col, bold=True,
           align=R, anchor=M)
        tb(s, vx + vzone - 42, yy, 36, h, unit, size=15, color=col, bold=True,
           align=L, anchor=M)
        if i < 3:
            cxx = fx + fw / 2
            add_icon(s, "chevron-right", cxx - 11, yy + h - 2, 22, GRAY2, stroke=2.2)
    # right: why it compounds
    rx = MARGIN + fw + 48; rw = 1920 - MARGIN - rx
    card(s, rx, y, rw, 458, fill=SOFT, line_c=LINE)
    tb(s, rx + 30, y + 28, rw - 60, 24, "学校チャネルの強さ", size=16, color=INK, bold=True)
    pts = [
        ("route", "低CAC", "学校1校の導入で学生・企業をまとめて獲得。"),
        ("repeat", "高い定着", "校務に組み込まれ、毎年新入生が自動流入。"),
        ("trending-up", "面の拡大", "学生→企業→学校と需給が相互に呼び込む。"),
    ]
    for k, (ic, t, d) in enumerate(pts):
        yy = y + 78 + k * 118
        icon_tile(s, rx + 30, yy, 50, ic)
        tb(s, rx + 96, yy - 2, rw - 126, 24, t, size=16, color=ACCENT7, bold=True)
        tb(s, rx + 96, yy + 26, rw - 126, 60, d, size=13.5, color=INK2, line=1.45)
    tb(s, rx + 30, y + 458 - 56, rw - 60, 40,
       "18ヶ月で 50校 → 3,000社（無料）→ 180社（有料）。", size=14,
       color=INK, bold=True, anchor=M)
    footer(s, 12)

# ----------------------------------------------------------------------------
# SLIDE 13 — Network Effects
# ----------------------------------------------------------------------------
def s13_network(prs):
    s = new_slide(prs)
    header(s, "NETWORK EFFECTS", "増えるほど、強くなる。",
           "三者とデータが循環し、後発が崩せない堀になる。")
    # flywheel: 4 nodes in a circle with curved-ish arrows
    cx, cy = 640, 610
    R = 215
    nodes = [
        ("school", "学校が増える", -90),
        ("graduation-cap", "学生・作品が増える", 0),
        ("building-2", "企業が集まる", 90),
        ("badge-check", "採用実績が貯まる", 180),
    ]
    import math
    pos = []
    for ic, label, ang in nodes:
        a = math.radians(ang)
        nx, ny = cx + R * math.cos(a), cy + R * math.sin(a)
        pos.append((nx, ny, ic, label))
    # ring arrows
    for i in range(4):
        x1, y1 = pos[i][0], pos[i][1]
        x2, y2 = pos[(i + 1) % 4][0], pos[(i + 1) % 4][1]
        line(s, x1, y1, x2, y2, color=ACCENT_TINT2, w=3.0)
    # center
    rect(s, cx - 70, cy - 70, 140, 140, fill=ACCENT, radius=0.5,
         shape=MSO_SHAPE.OVAL, shadow=True)
    add_icon(s, "repeat", cx - 30, cy - 38, 60, WHITE, stroke=1.9)
    tb(s, cx - 70, cy + 24, 140, 24, "Flywheel", size=13, color="C9DBFE",
       align=Cn, bold=True)
    for nx, ny, ic, label in pos:
        rect(s, nx - 56, ny - 56, 112, 112, fill=WHITE, line_c=ACCENT,
             line_w=1.6, radius=0.5, shape=MSO_SHAPE.OVAL, shadow=True)
        add_icon(s, ic, nx - 23, ny - 30, 46, ACCENT, stroke=1.8)
        # label clears the 56px node radius
        if abs(nx - cx) > R - 10:      # side nodes
            ly = ny + 64
        elif ny > cy:                  # bottom
            ly = ny + 66
        else:                          # top
            ly = ny - 84
        tb(s, nx - 150, ly, 300, 24, label, size=13.5, color=INK, bold=True, align=Cn)
    # right: data moat
    rx = 1140; rw = 1920 - MARGIN - rx; y = 350
    card(s, rx, y, rw, 470)
    tb(s, rx + 32, y + 30, rw - 64, 24, "データが、堀になる。", size=18, color=INK, bold=True)
    line(s, rx + 32, y + 70, rx + rw - 32, y + 70, color=LINE, w=1.0)
    moat = [
        ("image", "作品・スキルデータ", "学生の実力を示す独自データが蓄積。"),
        ("badge-check", "採用実績データ", "どのスキルが採用に繋がるかを学習。"),
        ("sparkles", "AIマッチング精度", "データが増えるほどマッチングが向上。"),
        ("school", "学校との関係資産", "毎年の新入生が自動で供給される。"),
    ]
    for k, (ic, t, d) in enumerate(moat):
        yy = y + 92 + k * 92
        icon_tile(s, rx + 32, yy, 46, ic)
        tb(s, rx + 92, yy - 2, rw - 120, 24, t, size=15, color=ACCENT7, bold=True)
        tb(s, rx + 92, yy + 24, rw - 120, 44, d, size=13, color=GRAY, line=1.35)
    footer(s, 13)

# ----------------------------------------------------------------------------
# SLIDE 14 — Traction
# ----------------------------------------------------------------------------
def s14_traction(prs):
    s = new_slide(prs)
    header(s, "TRACTION", "プロダクトは完成。市場は始動している。",
           "MVPとデモは完成済み。2026年7月の正式リリースに向け面が揃った。")
    stats = [
        ("badge-check", "MVP完成", "デモ提供中", "プロダクトは構築済み"),
        ("calendar", "2026年7月", "正式リリース", "ローンチ準備が完了"),
        ("school", "5校", "専門学校が始動", "初期パートナー校を確保"),
        ("building-2", "10社", "導入可能企業", "受け入れ企業が待機"),
    ]
    n = 4; gap = 26
    w = (CW - gap * (n - 1)) / n; y = 352; h = 230
    for i, (ic, v, l, sub) in enumerate(stats):
        x = MARGIN + i * (w + gap)
        card(s, x, y, w, h)
        icon_tile(s, x + 28, y + 28, 52, ic)
        tb(s, x + 28, y + 102, w - 56, 52, v, size=30, color=INK, bold=True, line=1.0)
        tb(s, x + 28, y + 168, w - 56, 24, l, size=15, color=ACCENT7, bold=True)
        tb(s, x + 28, y + 196, w - 56, 24, sub, size=12, color=GRAY)
    # milestone timeline
    ty = 640
    card(s, MARGIN, ty, CW, 196)
    tb(s, MARGIN + 30, ty + 24, CW - 60, 24, "マイルストーン", size=15, color=INK, bold=True)
    ms = [("構築", "高校市場で\n学校営業を実証"), ("開発", "MVP・デモ\n完成"),
          ("始動", "5校・10社\nオンボード"), ("ローンチ", "2026年7月\n正式リリース"),
          ("拡大", "50校・180社\n（18ヶ月）")]
    lx = MARGIN + 60; rxx = MARGIN + CW - 60
    base = ty + 118
    line(s, lx, base, rxx, base, color=LINE, w=2.0)
    seg = (rxx - lx) / (len(ms) - 1)
    for i, (t, d) in enumerate(ms):
        mx = lx + i * seg
        done = i < 4
        rect(s, mx - 11, base - 11, 22, 22, fill=ACCENT if done else WHITE,
             line_c=ACCENT, line_w=2.0, radius=0.5, shape=MSO_SHAPE.OVAL)
        if done:
            add_icon(s, "check", mx - 7, base - 7, 14, WHITE, stroke=3.0)
        tb(s, mx - 80, base - 56, 160, 24, t, size=13, color=ACCENT7 if done else GRAY,
           bold=True, align=Cn)
        tb(s, mx - 90, base + 18, 180, 50, d, size=11.5, color=GRAY, align=Cn, line=1.25)
    footer(s, 14)

# ----------------------------------------------------------------------------
# SLIDE 15 — Roadmap
# ----------------------------------------------------------------------------
def s15_roadmap(prs):
    s = new_slide(prs)
    header(s, "ROADMAP", "18ヶ月で、面を取り切る。",
           "正式リリースから、専門学校50校・有料企業180社へ。")
    # phase rollout bands (4 phases over timeline)
    y = 350
    timeline = ["2026 H2", "2027 H1", "2027 H2", "2028 H1"]
    tlx = MARGIN + 230; tlw = CW - 230
    seg = tlw / len(timeline)
    for i, tlbl in enumerate(timeline):
        tb(s, tlx + i * seg, y, seg, 22, tlbl, size=12.5, color=GRAY, bold=True, align=Cn)
    bands = [
        ("企業SaaS", 0, 4, ACCENT),
        ("スカウト", 1, 4, "3B7BF0"),
        ("採用課金", 2, 4, "60A5FA"),
        ("AI機能", 3, 4, "93C5FD"),
    ]
    by = y + 36; bh = 56; bgap = 16
    tb(s, MARGIN, by + 6, 210, 24, "プロダクト", size=12, color=GRAY2, bold=True, spc=1.0)
    for k, (name, start, end, col) in enumerate(bands):
        yy = by + k * (bh + bgap)
        rect(s, MARGIN, yy, 210, bh, fill=SOFT, line_c=LINE, radius=0.16,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb(s, MARGIN + 18, yy, 190, bh, name, size=14, color=INK, bold=True, anchor=M)
        bx = tlx + start * seg
        bw = (end - start) * seg - 12
        rect(s, bx, yy + 8, bw, bh - 16, fill=col, radius=0.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_icon(s, "arrow-right", bx + bw - 30, yy + bh / 2 - 11, 22, WHITE, stroke=2.2)
    # KPI targets row
    ky = by + 4 * (bh + bgap) + 18
    tb(s, MARGIN, ky, CW, 24, "18ヶ月後のKPI", size=15, color=INK, bold=True)
    kpis = [("school", "専門学校", "50校"), ("building-2", "無料企業", "3,000社"),
            ("credit-card", "有料企業", "180社"), ("graduation-cap", "学生", "数万名")]
    kw = (CW - 3 * 20) / 4; kyy = ky + 34
    for i, (ic, l, v) in enumerate(kpis):
        kx = MARGIN + i * (kw + 20)
        card(s, kx, kyy, kw, 116)
        icon_tile(s, kx + 24, kyy + 24, 46, ic)
        tb(s, kx + 86, kyy + 24, kw - 110, 20, l, size=13, color=GRAY, bold=True)
        tb(s, kx + 86, kyy + 46, kw - 110, 48, v, size=30, color=ACCENT7, bold=True)
    footer(s, 15)

# ----------------------------------------------------------------------------
# SLIDE 16 — Financial Plan
# ----------------------------------------------------------------------------
def s16_financial(prs):
    s = new_slide(prs)
    header(s, "FINANCIAL PLAN", "3年でARR ¥7.2億へ。",
           "学校チャネルの高い資本効率で、健全にスケールする。")
    # left: ARR growth bar chart
    lx = MARGIN; lw = 760; y = 348
    card(s, lx, y, lw, 372)
    tb(s, lx + 30, y + 26, lw - 60, 22, "ARR 成長", size=15, color=INK, bold=True)
    tb(s, lx + 30, y + 50, lw - 60, 20, "単位：百万円", size=11, color=GRAY2)
    bars = [("Y1", 65, "¥0.65億"), ("Y2", 260, "¥2.6億"), ("Y3", 720, "¥7.2億")]
    maxv = 720
    base = y + 300; chart_h = 190
    bw = 130; gap = (lw - 120 - bw * 3) / 2; x0 = lx + 70
    for i, (lab, v, vl) in enumerate(bars):
        bx = x0 + i * (bw + gap)
        bh = chart_h * (v / maxv)
        col = ACCENT if i == 2 else (ACCENT_TINT2 if i == 0 else "3B7BF0")
        rect(s, bx, base - bh, bw, bh, fill=col, radius=0.10,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb(s, bx - 20, base - bh - 34, bw + 40, 28, vl, size=16, color=INK,
           bold=True, align=Cn)
        tb(s, bx - 20, base + 12, bw + 40, 22, lab, size=13, color=GRAY,
           bold=True, align=Cn)
    line(s, x0 - 20, base, lx + lw - 40, base, color=LINE, w=1.2)
    # right: unit economics grid
    rx = lx + lw + 40; rw = 1920 - MARGIN - rx
    ue = [
        ("CAC", "¥12万", "学校チャネルで低く抑制"),
        ("LTV", "¥79万", "粗利80%・継続前提"),
        ("LTV / CAC", "6.6x", "健全な3x基準を大きく超過"),
        ("Payback", "5ヶ月", "回収が早く再投資可能"),
        ("粗利率", "80%超", "SaaS中心の高い利益率"),
        ("ARPA", "¥3万/月", "段階課金で拡張余地"),
    ]
    cols = 2; gap2 = 16
    cw = (rw - gap2) / 2; ch = 116; vg = 12
    for i, (l, v, d) in enumerate(ue):
        r, c = divmod(i, cols)
        cx = rx + c * (cw + gap2); cy = y + r * (ch + vg)
        card(s, cx, cy, cw, ch)
        tb(s, cx + 24, cy + 20, cw - 48, 20, l, size=12.5, color=GRAY, bold=True, spc=0.5)
        tb(s, cx + 24, cy + 44, cw - 48, 34, v, size=24, color=ACCENT7, bold=True, line=1.0)
        tb(s, cx + 24, cy + 90, cw - 48, 20, d, size=11, color=GRAY)
    # bottom funding/runway strip
    fy = 760
    seg = [("調達総額", "¥7,000万", "banknote"),
           ("内訳", "Equity ¥40M ＋ Debt ¥30M", "layers"),
           ("月次バーン", "約¥390万", "gauge"),
           ("Runway", "18ヶ月", "clock")]
    sw = (CW - 3 * 16) / 4
    for i, (l, v, ic) in enumerate(seg):
        sx = MARGIN + i * (sw + 16)
        rect(s, sx, fy, sw, 74, fill=INK if i == 0 else SOFT,
             line_c=None if i == 0 else LINE, radius=0.14,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_icon(s, ic, sx + 20, fy + 23, 28, WHITE if i == 0 else ACCENT7, stroke=1.9)
        tb(s, sx + 62, fy + 14, sw - 80, 20, l, size=11.5,
           color="C9DBFE" if i == 0 else GRAY, bold=True)
        tb(s, sx + 62, fy + 36, sw - 80, 28, v, size=16 if i != 1 else 12,
           color=WHITE if i == 0 else INK, bold=True, anchor=M)
    source(s, "数値は事業計画に基づく想定値（前提：ARPA ¥3万/月・粗利80%・月次解約率〜3%）。")
    footer(s, 16)

# ----------------------------------------------------------------------------
# SLIDE 17 — Team
# ----------------------------------------------------------------------------
def s17_team(prs):
    s = new_slide(prs)
    header(s, "TEAM", "学校営業を、現場で作ってきたチーム。",
           "プロダクト・営業・開発の三位一体で実行する。")
    members = [
        ("深川 練", "Founder", "学校営業モデルの設計者。高校市場で導入・運用網を構築。"),
        ("渡部", "CEO", "事業全体を統括。営業・パートナーシップを推進。"),
        ("片倉", "Engineer", "プロダクト開発をリード。MVP・デモを構築。"),
    ]
    n = 3; gap = 30
    w = (CW - gap * (n - 1)) / n; y = 352; h = 372
    for i, (name, role, desc) in enumerate(members):
        x = MARGIN + i * (w + gap)
        card(s, x, y, w, h)
        # photo placeholder
        ph = 120
        rect(s, x + w / 2 - ph / 2, y + 36, ph, ph, fill=SOFT, line_c=LINE,
             radius=0.5, shape=MSO_SHAPE.OVAL)
        add_icon(s, "user-round", x + w / 2 - 28, y + 36 + 28, 56, GRAY2, stroke=1.6)
        tb(s, x + 24, y + 176, w - 48, 36, name, size=22, color=INK, bold=True,
           align=Cn, line=1.0)
        pill(s, x + w / 2 - 80, y + 226, 160, 32, role, fill=ACCENT_TINT,
             color=ACCENT7, size=12.5, latin=EN)
        line(s, x + 40, y + 278, x + w - 40, y + 278, color=LINE, w=1.0)
        tb(s, x + 30, y + 296, w - 60, 60, desc, size=13.5, color=GRAY,
           align=Cn, line=1.45)
    # why this team
    by = 754
    rect(s, MARGIN, by, CW, 76, fill=ACCENT_TINT, radius=0.12,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_icon(s, "route", MARGIN + 28, by + 22, 32, ACCENT7, stroke=2.0)
    tb(s, MARGIN + 78, by, CW - 110, 76,
       "競争優位の源泉である「学校営業モデル」を、実際に高校市場で構築・実証してきた当事者チーム。",
       size=16, color=ACCENT7, bold=True, anchor=M)
    footer(s, 17)

# ----------------------------------------------------------------------------
# SLIDE 18 — Vision
# ----------------------------------------------------------------------------
def s18_vision(prs):
    s = new_slide(prs)
    # full ink background hero
    rect(s, 0, 0, 1920, 1080, fill=INK)
    rect(s, MARGIN, 150, 30, 30, fill=ACCENT, radius=0.28,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, MARGIN + 44, 150, 700, 30, "USTORY next", size=16, color=WHITE, bold=True)
    tb(s, MARGIN, 248, 1200, 60, "VISION", size=14, color="93C5FD", bold=True, spc=3.0)
    tb(s, MARGIN, 286, 1560, 220, "Build Japan's\nCareer Infrastructure.",
       size=60, color=WHITE, bold=True, line=1.08)
    tb(s, MARGIN, 556, 1300, 40,
       "才能が、学歴や履歴書に埋もれない社会をつくる。", size=21, color="D1D5DB", line=1.4)
    # mvp -> expansion arc
    arc = [("専門学校", "school"), ("全教育機関", "graduation-cap"),
           ("日本のキャリアインフラ", "globe")]
    y = 626; x0 = MARGIN; bw = 366; gap = 72
    for i, (t, ic) in enumerate(arc):
        bx = x0 + i * (bw + gap)
        accent = i == 2
        rect(s, bx, y, bw, 96, fill=ACCENT if accent else "1C2533",
             line_c=None if accent else "2C3647", radius=0.14,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_icon(s, ic, bx + 26, y + 28, 40, WHITE if accent else "93C5FD", stroke=1.8)
        tb(s, bx + 82, y, bw - 96, 96, t, size=17 if len(t) < 9 else 14.5,
           color=WHITE, bold=True, anchor=M)
        if i < 2:
            add_icon(s, "arrow-right", bx + bw + 20, y + 33, 30, "6B7280", stroke=2.2)
    # purpose + ask
    py = 778
    line(s, MARGIN, py, 1920 - MARGIN, py, color="2C3647", w=1.0)
    tb(s, MARGIN, py + 26, 1000, 30, "PURPOSE", size=13, color="93C5FD", bold=True, spc=2.0)
    tb(s, MARGIN, py + 54, 1120, 40,
       "Talent should be measured by skill, not resumes.", size=21, color=WHITE,
       bold=True, latin=EN)
    # ask box
    aw = 540
    ax = 1920 - MARGIN - aw
    rect(s, ax, py + 24, aw, 170, fill="1C2533", line_c="2C3647", radius=0.08,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, ax + 32, py + 40, aw - 64, 22, "THE ASK  ·  SEED ROUND", size=12,
       color="93C5FD", bold=True, spc=1.5)
    tb(s, ax + 32, py + 64, aw - 64, 46, "¥70,000,000", size=33, color=WHITE,
       bold=True, latin=EN, line=1.0)
    line(s, ax + 32, py + 124, ax + aw - 32, py + 124, color="2C3647", w=1.0)
    tb(s, ax + 32, py + 138, aw - 64, 22,
       "Equity ¥40M ＋ Debt ¥30M ／ 18ヶ月", size=13, color="D1D5DB")
    tb(s, MARGIN, 1024, 600, 22, "USTORY next", size=11, color="6B7280", bold=True)
    tb(s, 1920 - MARGIN - 300, 1024, 300, 22, "18 / 18", size=11, color="6B7280", align=R)

# ----------------------------------------------------------------------------
# Speaker notes (話す内容)
# ----------------------------------------------------------------------------
NOTES = [
 # 1 Cover
 "USTORY nextの深川です。本日は『日本のキャリアインフラを作る』という私たちの構想をお話しします。"
 "私たちは才能が学歴や履歴書に埋もれない社会を目指し、専門学校を起点に学校・学生・企業をつなぐ基盤を作っています。"
 "今回はシード7,000万円の調達についてご相談させてください。",
 # 2 Executive Summary
 "まず全体像です。事業は専門学校を起点にしたキャリアインフラ。市場はSAMで約2,851億円、年20〜30%成長。"
 "競争優位は高校市場で実証済みの『学校営業モデル』というDistribution。プロダクトは完成済みで、5校・10社と2026年7月に始動します。"
 "一言で言えば、専門スキルを持つ55万人の才能を、履歴書ではなくポートフォリオで評価される世界を作ります。",
 # 3 Investment Highlights
 "投資のポイントは6つ。成長市場、模倣困難なDistribution、強いネットワーク効果、実証済みの初期実績、高い資本効率、そして実行力あるチームです。"
 "特に強調したいのは、私たちの優位がAIではなく『学校への流通網』である点と、LTV/CAC 6.6倍という資本効率です。",
 # 4 Why Now
 "なぜ今か。3つの追い風が重なっています。新卒採用単価は1年で45万から56.8万へと高騰。"
 "HR Techクラウド市場はCAGR31.8%で3,200億円規模へ。そして2,676校ある専門学校はDXがほぼ手つかずで、スキル人材の流通網が存在しません。"
 "コストの問題、技術の追い風、そして空白市場。今こそ参入の好機です。",
 # 5 Problem
 "課題は、才能が履歴書に埋もれていることです。学生はスキルや作品があっても履歴書でしか評価されない。"
 "企業は専門スキル人材に出会えず採用単価だけが上がる。学校は就職支援が紙とExcelのアナログのまま。"
 "三者が断絶し、誰もこの溝を埋められていません。",
 # 6 Existing Solutions
 "既存サービスはこの三者の一辺ずつしか担えていません。求人媒体やスカウトは学生と企業をつなぐが学校とは連携しない。"
 "学校管理SaaSは校務を効率化するが採用と接続しない。ポートフォリオは見せる場止まり。"
 "結果、市場は断片化し、三者を同時につなぐキャリアインフラの空白が残っています。",
 # 7 Solution
 "私たちはこの三者を1つにつなぎます。学校を起点に、学生はスキルとポートフォリオを公開し、企業は検索・AI・スカウトから採用管理まで一気通貫で行える。"
 "学校は求人・説明会・学生管理をクラウド化。三者の体験が1つの基盤の上で循環します。",
 # 8 Product
 "プロダクトは3つの体験で構成されます。学生向けは作品・動画・スキルを載せるポートフォリオ。"
 "企業向けはスキル検索・AIマッチング・スカウトから採用管理まで。学校向けは求人・説明会・学生管理のクラウド。"
 "それぞれ独立した価値を持ちながら、1つのデータ基盤でつながっています。",
 # 9 Why USTORY
 "正直に申し上げると、私たちの競争優位はAIではありません。『学校営業モデル』というDistributionです。"
 "プロダクトは模倣されますが、学校との信頼関係と運用網は模倣できない。"
 "私たちは既に高校市場で学校営業・導入・CS・運用のモデルを構築・実証済みで、それをそのまま2,676校の専門学校市場へ展開します。",
 # 10 Market
 "市場規模です。最終的に挑む人材ビジネス市場全体はTAM10.2兆円。"
 "直接売り込むSAMは新卒採用支援とHR Techクラウドで2,851億円、年20〜30%成長。"
 "まず狙うSOMは、専門卒採用に動く約1,400億円の10%、5年で140億円。2,676校・55万人・採用単価56.8万円という確かな数字に基づいています。",
 # 11 Business Model
 "収益は4段階で積み上げます。まず企業SaaSで安定したMRRの土台を作り、次にスカウトの従量課金、"
 "そして採用課金の成果報酬で高単価レイヤーを上積み。最後にAIで単価と継続率をさらに高める。"
 "学校チャネルで企業を低CACで獲得し、ARPAを段階的に拡張する積み上げ型のモデルです。",
 # 12 Go To Market
 "GTMの肝は、学校を入口に企業を面で取ることです。学校1校の導入が数百人の学生と数十社の企業を連れてくる。"
 "学校チャネルは低CAC・高定着で、毎年新入生が自動的に流入します。18ヶ月で50校、無料企業3,000社、有料企業180社を目指します。",
 # 13 Network Effects
 "これは強いネットワーク効果を持つ事業です。学校が増えれば学生と作品が増え、企業が集まり、採用実績が貯まり、また学校が増える。"
 "さらに作品・スキル・採用実績のデータが蓄積し、AIマッチングの精度が上がる。増えるほど強くなり、後発が崩せない堀になります。",
 # 14 Traction
 "現在地です。プロダクトとデモは完成済み。2026年7月の正式リリースに向け、専門学校5校・導入可能企業10社という面が既に揃っています。"
 "高校市場で学校営業を実証し、開発を終え、いよいよローンチ。18ヶ月で50校・180社まで一気に拡大します。",
 # 15 Roadmap
 "ロードマップです。正式リリースから企業SaaS、スカウト、採用課金、AI機能と段階的に展開。"
 "18ヶ月後には専門学校50校、無料企業3,000社、有料企業180社、学生数万名というKPIを達成し、シリーズAにつなげます。",
 # 16 Financial Plan
 "財務計画です。3年でARRは0.65億→2.6億→7.2億円へ。学校チャネルでCACを12万円に抑え、LTVは79万円、LTV/CAC6.6倍、Payback5ヶ月、粗利80%超。"
 "今回の調達はEquity4,000万・Debt3,000万の計7,000万円で、月次バーン約390万円、18ヶ月のRunwayを確保します。",
 # 17 Team
 "チームです。私たちは競争優位の源泉である学校営業モデルを、実際に高校市場で構築・実証してきた当事者です。"
 "設計者の深川、事業を統括するCEO渡部、プロダクトをリードするエンジニア片倉。プロダクト・営業・開発の三位一体で実行します。",
 # 18 Vision
 "最後にビジョンです。私たちは専門学校から始め、全教育機関、そして日本のキャリアインフラへと広げていきます。"
 "才能はスキルで測られるべきで、履歴書で測られるべきではない。この実現のため、シード7,000万円の調達をお願いします。ぜひご一緒させてください。",
]

def set_notes(prs):
    for i, slide in enumerate(prs.slides):
        if i < len(NOTES):
            slide.notes_slide.notes_text_frame.text = NOTES[i]

# ----------------------------------------------------------------------------
# Build
# ----------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for fn in [s01_cover, s02_exec, s03_highlights, s04_whynow, s05_problem,
               s06_existing, s07_solution, s08_product, s09_why, s10_market,
               s11_model, s12_gtm, s13_network, s14_traction, s15_roadmap,
               s16_financial, s17_team, s18_vision]:
        fn(prs)
    set_notes(prs)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
